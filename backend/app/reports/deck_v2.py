"""
Enhanced PowerPoint deck generator with "Show Your Work" transparency.

This generator creates consultant-grade decks that:
1. Show the actual calculations behind each recommendation
2. Cite specific evidence from the data
3. Make assumptions explicit
4. Visualize confidence levels
5. Include actionable next steps with specificity
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor  # Correct import: RGBColor (capital RGB)
from pptx.enum.dml import MSO_THEME_COLOR
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from io import BytesIO


# Color palette - professional consulting style
class Colors:
    # Primary
    DARK_BLUE = RgbColor(0x1a, 0x36, 0x5d)      # Headers, titles
    MEDIUM_BLUE = RgbColor(0x2c, 0x5a, 0x8a)    # Accents
    LIGHT_BLUE = RgbColor(0xe8, 0xf4, 0xfc)     # Backgrounds

    # Confidence indicators
    HIGH_CONFIDENCE = RgbColor(0x2e, 0x7d, 0x32)    # Green
    MEDIUM_CONFIDENCE = RgbColor(0xf5, 0x9e, 0x0b)  # Amber
    LOW_CONFIDENCE = RgbColor(0xc6, 0x28, 0x28)     # Red

    # Neutral
    DARK_GRAY = RgbColor(0x33, 0x33, 0x33)
    MEDIUM_GRAY = RgbColor(0x66, 0x66, 0x66)
    LIGHT_GRAY = RgbColor(0xf5, 0xf5, 0xf5)
    WHITE = RgbColor(0xff, 0xff, 0xff)


class EnhancedDeckGenerator:
    """
    Generate PowerPoint presentations with full transparency on calculations.

    Philosophy: Every number should be traceable. Every recommendation
    should show the math. Every assumption should be explicit.
    """

    def __init__(self):
        self.prs = None
        self.slide_width = Inches(13.333)  # Widescreen 16:9
        self.slide_height = Inches(7.5)

    def generate(
        self,
        company_name: str,
        mode_info: Dict[str, Any],
        analytics_facts: List[Dict[str, Any]],
        initiatives: List[Dict[str, Any]],
        vertical_name: str,
        output_path: str,
        run_context: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Generate enhanced PowerPoint deck.

        Args:
            company_name: Company name
            mode_info: Operating mode information (mode, confidence, months_available)
            analytics_facts: List of computed analytics facts
            initiatives: Selected and sized initiatives with specificity drafts
            vertical_name: Vertical name (e.g., "Restaurant Operations")
            output_path: Output file path
            run_context: Optional run context from intake questions

        Returns:
            Path to generated file
        """
        self.prs = Presentation()
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height

        # Build facts lookup for easy reference
        facts_map = {fact['evidence_key']: fact for fact in analytics_facts}

        # === SLIDE 1: Title ===
        self._add_title_slide(company_name, vertical_name, mode_info)

        # === SLIDE 2: Executive Summary ===
        self._add_executive_summary(mode_info, analytics_facts, initiatives)

        # === SLIDE 3: Your Numbers at a Glance (with chart) ===
        self._add_metrics_dashboard(analytics_facts, mode_info)

        # === SLIDE 4: How We Analyzed Your Data ===
        self._add_methodology_slide(mode_info, analytics_facts)

        # === SLIDES 5-7: Top 3 Initiative Deep Dives ===
        for initiative in initiatives[:3]:
            self._add_initiative_deep_dive(initiative, facts_map, mode_info)

        # === SLIDE 8: All Initiatives Ranked ===
        self._add_initiatives_summary(initiatives)

        # === SLIDE 9: Our Assumptions ===
        self._add_assumptions_slide(initiatives, mode_info)

        # === SLIDE 10: What Would Make This Better ===
        self._add_data_gaps_slide(initiatives, mode_info)

        # Save
        self.prs.save(output_path)
        return output_path

    # =========================================================================
    # SLIDE BUILDERS
    # =========================================================================

    def _add_title_slide(
        self,
        company_name: str,
        vertical_name: str,
        mode_info: Dict[str, Any]
    ):
        """Slide 1: Clean title slide."""
        slide = self._add_blank_slide()

        # Company name - large, centered
        self._add_text_box(
            slide,
            text=company_name,
            left=Inches(1),
            top=Inches(2.5),
            width=Inches(11.333),
            height=Inches(1.2),
            font_size=54,
            font_bold=True,
            font_color=Colors.DARK_BLUE,
            alignment=PP_ALIGN.CENTER
        )

        # Subtitle
        self._add_text_box(
            slide,
            text="Operational Diagnostic Report",
            left=Inches(1),
            top=Inches(3.7),
            width=Inches(11.333),
            height=Inches(0.6),
            font_size=28,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Meta info
        confidence = mode_info.get('confidence', 0)
        confidence_pct = f"{confidence * 100:.0f}%" if isinstance(confidence, float) else str(confidence)
        meta_text = f"{vertical_name}  •  {mode_info.get('months_available', 0)} months analyzed  •  {confidence_pct} confidence"

        self._add_text_box(
            slide,
            text=meta_text,
            left=Inches(1),
            top=Inches(6.5),
            width=Inches(11.333),
            height=Inches(0.4),
            font_size=14,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Date
        self._add_text_box(
            slide,
            text=datetime.now().strftime("%B %Y"),
            left=Inches(1),
            top=Inches(6.9),
            width=Inches(11.333),
            height=Inches(0.3),
            font_size=12,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    def _add_executive_summary(
        self,
        mode_info: Dict[str, Any],
        analytics_facts: List[Dict[str, Any]],
        initiatives: List[Dict[str, Any]]
    ):
        """Slide 2: Executive summary with key findings."""
        slide = self._add_blank_slide()

        # Title
        self._add_slide_title(slide, "Executive Summary")

        # Calculate total potential impact
        total_low = sum(init.get('impact_low', 0) for init in initiatives)
        total_high = sum(init.get('impact_high', 0) for init in initiatives)

        # Big number - total opportunity
        self._add_text_box(
            slide,
            text=f"${total_low:,.0f} – ${total_high:,.0f}",
            left=Inches(1),
            top=Inches(1.8),
            width=Inches(11.333),
            height=Inches(0.9),
            font_size=48,
            font_bold=True,
            font_color=Colors.DARK_BLUE,
            alignment=PP_ALIGN.CENTER
        )

        self._add_text_box(
            slide,
            text="Estimated Annual Impact Opportunity",
            left=Inches(1),
            top=Inches(2.7),
            width=Inches(11.333),
            height=Inches(0.4),
            font_size=16,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Three columns: Mode | Top Initiative | Key Metric
        col_width = Inches(3.5)
        col_start = Inches(1)
        col_top = Inches(3.5)

        # Column 1: Operating Mode
        mode = mode_info.get('mode', 'DIRECTIONAL_MODE')
        mode_display = mode.replace('_MODE', '').replace('_', ' ').title()

        self._add_metric_card(
            slide,
            label="Analysis Mode",
            value=mode_display,
            detail=f"{mode_info.get('months_available', 0)} months of data",
            left=col_start,
            top=col_top,
            width=col_width
        )

        # Column 2: Top Initiative
        if initiatives:
            top_init = initiatives[0]
            self._add_metric_card(
                slide,
                label="Top Opportunity",
                value=top_init.get('title', 'N/A')[:30],
                detail=f"${top_init.get('impact_mid', 0):,.0f} potential",
                left=col_start + col_width + Inches(0.5),
                top=col_top,
                width=col_width
            )

        # Column 3: Key Metric (find most important fact)
        key_facts = ['revenue_avg_monthly', 'labor_pct', 'revenue_trend']
        key_fact = None
        for key in key_facts:
            for fact in analytics_facts:
                if fact.get('evidence_key') == key:
                    key_fact = fact
                    break
            if key_fact:
                break

        if key_fact:
            value = key_fact.get('value', key_fact.get('value_text', 'N/A'))
            if isinstance(value, float):
                if key_fact.get('unit') == 'percentage':
                    value_str = f"{value:.1f}%"
                elif key_fact.get('unit') == 'currency':
                    value_str = f"${value:,.0f}"
                else:
                    value_str = f"{value:,.1f}"
            else:
                value_str = str(value)

            self._add_metric_card(
                slide,
                label=key_fact.get('label', 'Key Metric'),
                value=value_str,
                detail=key_fact.get('evidence_key', ''),
                left=col_start + (col_width + Inches(0.5)) * 2,
                top=col_top,
                width=col_width
            )

        # Bottom: Confidence statement
        confidence = mode_info.get('confidence', 0)
        if confidence >= 0.7:
            conf_text = "High confidence analysis based on comprehensive data coverage."
            conf_color = Colors.HIGH_CONFIDENCE
        elif confidence >= 0.5:
            conf_text = "Moderate confidence analysis. Some data gaps exist."
            conf_color = Colors.MEDIUM_CONFIDENCE
        else:
            conf_text = "Directional analysis only. Additional data recommended for validation."
            conf_color = Colors.LOW_CONFIDENCE

        # Confidence indicator bar
        self._add_confidence_bar(slide, confidence, Inches(1), Inches(5.8), Inches(11.333))

        self._add_text_box(
            slide,
            text=conf_text,
            left=Inches(1),
            top=Inches(6.3),
            width=Inches(11.333),
            height=Inches(0.4),
            font_size=12,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    def _add_metrics_dashboard(
        self,
        analytics_facts: List[Dict[str, Any]],
        mode_info: Dict[str, Any]
    ):
        """Slide 3: Key metrics with visualization."""
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "Your Numbers at a Glance")

        # Left side: Key metrics list
        metrics_to_show = [
            'revenue_avg_monthly',
            'labor_pct',
            'cogs_pct',
            'revenue_trend',
            'labor_volatility'
        ]

        y_pos = Inches(1.8)
        for evidence_key in metrics_to_show:
            fact = next((f for f in analytics_facts if f.get('evidence_key') == evidence_key), None)
            if fact:
                self._add_metric_row(slide, fact, Inches(0.8), y_pos, Inches(5.5))
                y_pos += Inches(0.9)

        # Right side: Simple bar chart of key ratios
        ratio_facts = [f for f in analytics_facts if f.get('unit') == 'percentage'][:4]
        if ratio_facts:
            chart_image = self._create_ratio_chart(ratio_facts)
            if chart_image:
                slide.shapes.add_picture(
                    chart_image,
                    Inches(7),
                    Inches(1.8),
                    width=Inches(5.5),
                    height=Inches(4)
                )

    def _add_methodology_slide(
        self,
        mode_info: Dict[str, Any],
        analytics_facts: List[Dict[str, Any]]
    ):
        """Slide 4: How we analyzed the data - transparency on methodology."""
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "How We Analyzed Your Data")

        # Data coverage section
        months = mode_info.get('months_available', 0)
        mode = mode_info.get('mode', 'DIRECTIONAL_MODE')
        confidence = mode_info.get('confidence', 0)

        # Left column: What we had
        self._add_text_box(
            slide,
            text="DATA INPUTS",
            left=Inches(0.8),
            top=Inches(1.8),
            width=Inches(5.5),
            height=Inches(0.4),
            font_size=14,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        data_text = f"""• {months} months of financial data
• {len(analytics_facts)} metrics computed
• Operating mode: {mode.replace('_', ' ').title()}
• Overall confidence: {confidence * 100:.0f}%"""

        self._add_text_box(
            slide,
            text=data_text,
            left=Inches(0.8),
            top=Inches(2.3),
            width=Inches(5.5),
            height=Inches(2),
            font_size=14,
            font_color=Colors.DARK_GRAY
        )

        # Right column: How we calculated
        self._add_text_box(
            slide,
            text="OUR APPROACH",
            left=Inches(7),
            top=Inches(1.8),
            width=Inches(5.5),
            height=Inches(0.4),
            font_size=14,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        approach_text = """• All metrics computed deterministically (no AI guessing)
• Impact estimates based on your actual numbers
• Industry benchmarks used where data gaps exist
• Recommendations filtered by your constraints
• Confidence levels reflect data quality"""

        self._add_text_box(
            slide,
            text=approach_text,
            left=Inches(7),
            top=Inches(2.3),
            width=Inches(5.5),
            height=Inches(2),
            font_size=14,
            font_color=Colors.DARK_GRAY
        )

        # Bottom: Mode explanation
        mode_explanations = {
            'PNL_MODE': "Full P&L analysis enabled. High confidence in cost-based recommendations.",
            'OPS_MODE': "Operational analysis mode. Revenue and transaction patterns available.",
            'DIRECTIONAL_MODE': "Limited data available. Recommendations are directional guidance only."
        }

        self._add_text_box(
            slide,
            text=mode_explanations.get(mode, "Analysis complete."),
            left=Inches(0.8),
            top=Inches(5.5),
            width=Inches(11.733),
            height=Inches(0.5),
            font_size=14,
            font_color=Colors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER,
            italic=True
        )

    def _add_initiative_deep_dive(
        self,
        initiative: Dict[str, Any],
        facts_map: Dict[str, Dict[str, Any]],
        mode_info: Dict[str, Any]
    ):
        """
        Slide 5-7: Deep dive on a single initiative.

        THE KEY SLIDE - This is where we "show our work"
        """
        slide = self._add_blank_slide()

        rank = initiative.get('rank', '?')
        title = initiative.get('title', 'Initiative')

        # Title with rank badge
        self._add_slide_title(slide, f"#{rank}: {title}")

        # Get specificity draft
        spec = initiative.get('specificity_draft', {})

        # === LEFT COLUMN: The Recommendation ===
        left_col = Inches(0.8)
        col_width = Inches(5.8)

        # Impact range - THE BIG NUMBER
        impact_low = initiative.get('impact_low', 0)
        impact_high = initiative.get('impact_high', 0)

        self._add_text_box(
            slide,
            text=f"${impact_low:,.0f} – ${impact_high:,.0f}",
            left=left_col,
            top=Inches(1.6),
            width=col_width,
            height=Inches(0.7),
            font_size=36,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        self._add_text_box(
            slide,
            text="Estimated Annual Impact",
            left=left_col,
            top=Inches(2.2),
            width=col_width,
            height=Inches(0.3),
            font_size=12,
            font_color=Colors.MEDIUM_GRAY
        )

        # THE CALCULATION - Show the math
        self._add_text_box(
            slide,
            text="HOW WE CALCULATED THIS",
            left=left_col,
            top=Inches(2.8),
            width=col_width,
            height=Inches(0.35),
            font_size=12,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        calc_text = self._build_calculation_explanation(initiative, facts_map)

        self._add_text_box(
            slide,
            text=calc_text,
            left=left_col,
            top=Inches(3.15),
            width=col_width,
            height=Inches(1.8),
            font_size=11,
            font_color=Colors.DARK_GRAY
        )

        # === RIGHT COLUMN: Action Plan ===
        right_col = Inches(7)

        # What to do
        self._add_text_box(
            slide,
            text="WHAT TO DO",
            left=right_col,
            top=Inches(1.6),
            width=col_width,
            height=Inches(0.35),
            font_size=12,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        what_text = spec.get('what', initiative.get('description', ''))
        where_text = spec.get('where', '')
        how_much_text = spec.get('how_much', '')

        action_text = what_text
        if where_text and 'NOT APPLICABLE' not in where_text:
            action_text += f"\n\nWhere: {where_text}"
        if how_much_text:
            action_text += f"\n\nTarget: {how_much_text}"

        self._add_text_box(
            slide,
            text=action_text,
            left=right_col,
            top=Inches(1.95),
            width=col_width,
            height=Inches(1.6),
            font_size=11,
            font_color=Colors.DARK_GRAY
        )

        # Next steps
        next_steps = spec.get('next_steps', [])
        if next_steps:
            self._add_text_box(
                slide,
                text="NEXT STEPS",
                left=right_col,
                top=Inches(3.6),
                width=col_width,
                height=Inches(0.35),
                font_size=12,
                font_bold=True,
                font_color=Colors.DARK_BLUE
            )

            steps_text = "\n".join([f"☐ {step}" for step in next_steps[:5]])

            self._add_text_box(
                slide,
                text=steps_text,
                left=right_col,
                top=Inches(3.95),
                width=col_width,
                height=Inches(1.5),
                font_size=10,
                font_color=Colors.DARK_GRAY
            )

        # === BOTTOM: Confidence & Assumptions ===
        confidence = spec.get('confidence', 'MEDIUM')
        specificity = spec.get('specificity_level', 'DIRECTIONAL')
        assumptions = spec.get('assumptions', initiative.get('assumptions', []))

        # Confidence badge
        conf_colors = {
            'HIGH': Colors.HIGH_CONFIDENCE,
            'MEDIUM': Colors.MEDIUM_CONFIDENCE,
            'LOW': Colors.LOW_CONFIDENCE,
            'N/A': Colors.MEDIUM_GRAY
        }

        self._add_text_box(
            slide,
            text=f"Confidence: {confidence}  |  Specificity: {specificity}",
            left=left_col,
            top=Inches(5.4),
            width=Inches(5),
            height=Inches(0.35),
            font_size=11,
            font_bold=True,
            font_color=conf_colors.get(confidence, Colors.MEDIUM_GRAY)
        )

        # Key assumptions
        if assumptions:
            assumptions_text = "Key assumptions: " + " • ".join(assumptions[:3])
            self._add_text_box(
                slide,
                text=assumptions_text,
                left=left_col,
                top=Inches(5.8),
                width=Inches(11.733),
                height=Inches(0.6),
                font_size=10,
                font_color=Colors.MEDIUM_GRAY,
                italic=True
            )

    def _add_initiatives_summary(self, initiatives: List[Dict[str, Any]]):
        """Slide 8: All initiatives ranked in a table."""
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "All Recommended Initiatives")

        # Create table
        rows = min(len(initiatives), 8) + 1  # Header + up to 8 initiatives
        cols = 5

        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.6),
            Inches(12.333), Inches(4.5)
        ).table

        # Set column widths
        table.columns[0].width = Inches(0.6)   # Rank
        table.columns[1].width = Inches(4)     # Initiative
        table.columns[2].width = Inches(2.2)   # Category
        table.columns[3].width = Inches(2.8)   # Impact Range
        table.columns[4].width = Inches(2.733) # Confidence

        # Header row
        headers = ['#', 'Initiative', 'Category', 'Annual Impact', 'Confidence']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = Colors.DARK_BLUE
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = Colors.WHITE
            para.font.size = Pt(11)
            para.font.bold = True

        # Data rows
        for row_idx, init in enumerate(initiatives[:8]):
            row = row_idx + 1

            # Rank
            table.cell(row, 0).text = str(init.get('rank', row))

            # Title
            table.cell(row, 1).text = init.get('title', 'N/A')

            # Category
            table.cell(row, 2).text = init.get('category', 'N/A')

            # Impact
            low = init.get('impact_low', 0)
            high = init.get('impact_high', 0)
            table.cell(row, 3).text = f"${low:,.0f} – ${high:,.0f}"

            # Confidence
            spec = init.get('specificity_draft', {})
            conf = spec.get('confidence', 'MEDIUM')
            table.cell(row, 4).text = conf

            # Style data cells
            for col in range(cols):
                cell = table.cell(row, col)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = Colors.DARK_GRAY

                # Alternate row colors
                if row % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Colors.LIGHT_GRAY

    def _add_assumptions_slide(
        self,
        initiatives: List[Dict[str, Any]],
        mode_info: Dict[str, Any]
    ):
        """Slide 9: All assumptions made transparent."""
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "Our Assumptions")

        self._add_text_box(
            slide,
            text="These recommendations assume the following. Please validate against your knowledge of the business.",
            left=Inches(0.8),
            top=Inches(1.5),
            width=Inches(11.733),
            height=Inches(0.5),
            font_size=14,
            font_color=Colors.MEDIUM_GRAY,
            italic=True
        )

        # Collect all unique assumptions
        all_assumptions = set()

        # Mode-based assumptions
        mode = mode_info.get('mode', '')
        if mode == 'DIRECTIONAL_MODE':
            all_assumptions.add("Limited data - recommendations are directional guidance only")
        elif mode == 'OPS_MODE':
            all_assumptions.add("P&L data not available - cost estimates use industry benchmarks")

        # Initiative assumptions
        for init in initiatives:
            spec = init.get('specificity_draft', {})
            for assumption in spec.get('assumptions', []):
                all_assumptions.add(assumption)
            for assumption in init.get('assumptions', []):
                all_assumptions.add(assumption)

        # Display assumptions in two columns
        assumptions_list = list(all_assumptions)[:16]  # Max 16
        mid = len(assumptions_list) // 2 + len(assumptions_list) % 2

        left_assumptions = assumptions_list[:mid]
        right_assumptions = assumptions_list[mid:]

        if left_assumptions:
            left_text = "\n\n".join([f"• {a}" for a in left_assumptions])
            self._add_text_box(
                slide,
                text=left_text,
                left=Inches(0.8),
                top=Inches(2.2),
                width=Inches(5.8),
                height=Inches(4.5),
                font_size=11,
                font_color=Colors.DARK_GRAY
            )

        if right_assumptions:
            right_text = "\n\n".join([f"• {a}" for a in right_assumptions])
            self._add_text_box(
                slide,
                text=right_text,
                left=Inches(7),
                top=Inches(2.2),
                width=Inches(5.8),
                height=Inches(4.5),
                font_size=11,
                font_color=Colors.DARK_GRAY
            )

    def _add_data_gaps_slide(
        self,
        initiatives: List[Dict[str, Any]],
        mode_info: Dict[str, Any]
    ):
        """Slide 10: What additional data would improve this analysis."""
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "What Would Make This Better")

        self._add_text_box(
            slide,
            text="Providing the following data would increase our confidence and specificity:",
            left=Inches(0.8),
            top=Inches(1.5),
            width=Inches(11.733),
            height=Inches(0.5),
            font_size=14,
            font_color=Colors.MEDIUM_GRAY,
            italic=True
        )

        # Collect all data gaps
        all_gaps = set()

        # Mode-based gaps
        mode = mode_info.get('mode', '')
        if mode == 'DIRECTIONAL_MODE':
            all_gaps.add("At least 3 months of P&L data for full analysis")
        if mode in ['DIRECTIONAL_MODE', 'OPS_MODE']:
            all_gaps.add("Detailed expense breakdown by category")

        # Initiative-specific gaps
        for init in initiatives:
            spec = init.get('specificity_draft', {})
            for gap in spec.get('data_needed', []):
                all_gaps.add(gap)
            for gap in init.get('data_gaps', []):
                all_gaps.add(gap)

        # Display in two columns
        gaps_list = list(all_gaps)[:12]  # Max 12
        mid = len(gaps_list) // 2 + len(gaps_list) % 2

        left_gaps = gaps_list[:mid]
        right_gaps = gaps_list[mid:]

        if left_gaps:
            left_text = "\n\n".join([f"☐ {g}" for g in left_gaps])
            self._add_text_box(
                slide,
                text=left_text,
                left=Inches(0.8),
                top=Inches(2.2),
                width=Inches(5.8),
                height=Inches(3.5),
                font_size=12,
                font_color=Colors.DARK_GRAY
            )

        if right_gaps:
            right_text = "\n\n".join([f"☐ {g}" for g in right_gaps])
            self._add_text_box(
                slide,
                text=right_text,
                left=Inches(7),
                top=Inches(2.2),
                width=Inches(5.8),
                height=Inches(3.5),
                font_size=12,
                font_color=Colors.DARK_GRAY
            )

        # Bottom CTA
        self._add_text_box(
            slide,
            text="Ready to go deeper? Upload additional data for a more detailed analysis.",
            left=Inches(0.8),
            top=Inches(6.2),
            width=Inches(11.733),
            height=Inches(0.5),
            font_size=14,
            font_bold=True,
            font_color=Colors.DARK_BLUE,
            alignment=PP_ALIGN.CENTER
        )

    # =========================================================================
    # HELPER METHODS
    # =========================================================================

    def _add_blank_slide(self):
        """Add a blank slide."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(blank_layout)

    def _add_slide_title(self, slide, title: str):
        """Add standardized slide title."""
        self._add_text_box(
            slide,
            text=title,
            left=Inches(0.8),
            top=Inches(0.5),
            width=Inches(11.733),
            height=Inches(0.8),
            font_size=32,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

    def _add_text_box(
        self,
        slide,
        text: str,
        left,
        top,
        width,
        height,
        font_size: int = 14,
        font_bold: bool = False,
        font_color: RgbColor = None,
        alignment=PP_ALIGN.LEFT,
        italic: bool = False
    ):
        """Add a text box with styling."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = font_bold
        p.font.italic = italic
        p.alignment = alignment

        if font_color:
            p.font.color.rgb = font_color

        return txBox

    def _add_metric_card(
        self,
        slide,
        label: str,
        value: str,
        detail: str,
        left,
        top,
        width
    ):
        """Add a metric card with label, value, and detail."""
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.LIGHT_GRAY
        shape.line.fill.background()

        # Label
        self._add_text_box(
            slide,
            text=label.upper(),
            left=left + Inches(0.2),
            top=top + Inches(0.15),
            width=width - Inches(0.4),
            height=Inches(0.3),
            font_size=10,
            font_color=Colors.MEDIUM_GRAY
        )

        # Value
        self._add_text_box(
            slide,
            text=value,
            left=left + Inches(0.2),
            top=top + Inches(0.45),
            width=width - Inches(0.4),
            height=Inches(0.5),
            font_size=18,
            font_bold=True,
            font_color=Colors.DARK_BLUE
        )

        # Detail
        self._add_text_box(
            slide,
            text=detail,
            left=left + Inches(0.2),
            top=top + Inches(1),
            width=width - Inches(0.4),
            height=Inches(0.3),
            font_size=10,
            font_color=Colors.MEDIUM_GRAY
        )

    def _add_metric_row(
        self,
        slide,
        fact: Dict[str, Any],
        left,
        top,
        width
    ):
        """Add a single metric row."""
        label = fact.get('label', fact.get('evidence_key', 'Metric'))
        value = fact.get('value', fact.get('value_text', 'N/A'))
        unit = fact.get('unit', '')

        # Format value
        if isinstance(value, float):
            if unit == 'percentage':
                value_str = f"{value:.1f}%"
            elif unit == 'currency':
                value_str = f"${value:,.0f}"
            else:
                value_str = f"{value:,.2f}"
        else:
            value_str = str(value)

        # Label
        self._add_text_box(
            slide,
            text=label,
            left=left,
            top=top,
            width=Inches(3.5),
            height=Inches(0.4),
            font_size=14,
            font_color=Colors.DARK_GRAY
        )

        # Value
        self._add_text_box(
            slide,
            text=value_str,
            left=left + Inches(3.5),
            top=top,
            width=Inches(2),
            height=Inches(0.4),
            font_size=14,
            font_bold=True,
            font_color=Colors.DARK_BLUE,
            alignment=PP_ALIGN.RIGHT
        )

        # Evidence key (small)
        self._add_text_box(
            slide,
            text=fact.get('evidence_key', ''),
            left=left,
            top=top + Inches(0.35),
            width=width,
            height=Inches(0.25),
            font_size=9,
            font_color=Colors.MEDIUM_GRAY
        )

    def _add_confidence_bar(
        self,
        slide,
        confidence: float,
        left,
        top,
        width
    ):
        """Add a visual confidence bar."""
        bar_height = Inches(0.15)

        # Background bar
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, bar_height
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = Colors.LIGHT_GRAY
        bg_bar.line.fill.background()

        # Confidence fill
        fill_width = width * min(confidence, 1.0)
        if fill_width > Inches(0.1):  # Minimum visible width
            fill_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, fill_width, bar_height
            )

            if confidence >= 0.7:
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = Colors.HIGH_CONFIDENCE
            elif confidence >= 0.5:
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = Colors.MEDIUM_CONFIDENCE
            else:
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = Colors.LOW_CONFIDENCE

            fill_bar.line.fill.background()

    def _build_calculation_explanation(
        self,
        initiative: Dict[str, Any],
        facts_map: Dict[str, Dict[str, Any]]
    ) -> str:
        """
        Build a human-readable explanation of how we calculated the impact.

        THIS IS THE KEY "SHOW YOUR WORK" FUNCTION.
        """
        method = initiative.get('sizing_method', 'fixed_value')
        params = initiative.get('sizing_params', {})

        explanation_parts = []

        if method == 'percentage_of_revenue':
            revenue_fact = facts_map.get('revenue_avg_monthly')
            if revenue_fact:
                monthly = revenue_fact.get('value', 0)
                annual = monthly * 12
                low_pct = params.get('low', 0) * 100
                high_pct = params.get('high', 0) * 100

                explanation_parts.append(f"Your average monthly revenue: ${monthly:,.0f}")
                explanation_parts.append(f"Annualized: ${annual:,.0f}")
                explanation_parts.append(f"")
                explanation_parts.append(f"Impact range: {low_pct:.0f}% to {high_pct:.0f}% of annual revenue")
                explanation_parts.append(f"")
                explanation_parts.append(f"Low estimate: ${annual:,.0f} × {low_pct:.0f}% = ${annual * params.get('low', 0):,.0f}")
                explanation_parts.append(f"High estimate: ${annual:,.0f} × {high_pct:.0f}% = ${annual * params.get('high', 0):,.0f}")
            else:
                explanation_parts.append("Revenue data not available.")
                explanation_parts.append("Using industry benchmark estimates.")

        elif method == 'percentage_of_labor':
            labor_fact = facts_map.get('labor_avg_monthly')
            if labor_fact:
                monthly = labor_fact.get('value', 0)
                annual = monthly * 12
                low_pct = params.get('low', 0) * 100
                high_pct = params.get('high', 0) * 100

                explanation_parts.append(f"Your average monthly labor cost: ${monthly:,.0f}")
                explanation_parts.append(f"Annualized: ${annual:,.0f}")
                explanation_parts.append(f"")
                explanation_parts.append(f"Potential savings: {low_pct:.0f}% to {high_pct:.0f}% of labor costs")
                explanation_parts.append(f"")
                explanation_parts.append(f"Low estimate: ${annual:,.0f} × {low_pct:.0f}% = ${annual * params.get('low', 0):,.0f}")
                explanation_parts.append(f"High estimate: ${annual:,.0f} × {high_pct:.0f}% = ${annual * params.get('high', 0):,.0f}")
            else:
                explanation_parts.append("Labor cost data not available.")
                explanation_parts.append("Using industry benchmark estimates.")

        elif method == 'percentage_of_cogs':
            # Try to find COGS fact
            cogs_value = None
            for key, fact in facts_map.items():
                if 'cogs' in key.lower():
                    cogs_value = fact.get('value')
                    break

            if cogs_value:
                annual = cogs_value * 12
                low_pct = params.get('low', 0) * 100
                high_pct = params.get('high', 0) * 100

                explanation_parts.append(f"Your estimated monthly COGS: ${cogs_value:,.0f}")
                explanation_parts.append(f"Annualized: ${annual:,.0f}")
                explanation_parts.append(f"")
                explanation_parts.append(f"Potential savings: {low_pct:.0f}% to {high_pct:.0f}% of COGS")
            else:
                explanation_parts.append("COGS data not available.")
                explanation_parts.append("Using industry benchmark estimates.")

        elif method == 'fixed_value':
            low = params.get('low', 0)
            mid = params.get('mid', 0)
            high = params.get('high', 0)

            explanation_parts.append("Impact based on industry benchmarks")
            explanation_parts.append("for similar initiatives:")
            explanation_parts.append(f"")
            explanation_parts.append(f"Conservative: ${low:,.0f}")
            explanation_parts.append(f"Expected: ${mid:,.0f}")
            explanation_parts.append(f"Optimistic: ${high:,.0f}")

        else:
            explanation_parts.append("Calculation method: " + method)
            explanation_parts.append("Based on industry benchmarks and your data.")

        return "\n".join(explanation_parts)

    def _create_ratio_chart(
        self,
        ratio_facts: List[Dict[str, Any]]
    ) -> Optional[BytesIO]:
        """Create a horizontal bar chart of ratio metrics."""
        try:
            fig, ax = plt.subplots(figsize=(6, 4))

            labels = [f.get('label', f.get('evidence_key', ''))[:25] for f in ratio_facts]
            values = [f.get('value', 0) for f in ratio_facts]

            # Reverse for horizontal bar (top to bottom)
            labels = labels[::-1]
            values = values[::-1]

            bars = ax.barh(labels, values, color='#2c5a8a')

            # Add value labels
            for bar, value in zip(bars, values):
                ax.text(
                    bar.get_width() + 0.5,
                    bar.get_y() + bar.get_height() / 2,
                    f'{value:.1f}%',
                    va='center',
                    fontsize=10
                )

            ax.set_xlabel('Percentage')
            ax.set_xlim(0, max(values) * 1.3 if values else 100)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)

            plt.tight_layout()

            # Save to bytes
            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            plt.close(fig)

            return buf

        except Exception as e:
            return None
