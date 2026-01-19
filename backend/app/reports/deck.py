"""Generate PowerPoint deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Any
from datetime import datetime
import matplotlib.pyplot as plt
import pandas as pd
from io import BytesIO


class DeckGenerator:
    """Generate PowerPoint presentation."""
    
    def generate(
        self,
        company_name: str,
        mode_info: Dict[str, Any],
        analytics_facts: List[Dict[str, Any]],
        initiatives: List[Dict[str, Any]],
        vertical_name: str,
        output_path: str
    ):
        """
        Generate PowerPoint deck.
        
        Args:
            company_name: Company name
            mode_info: Operating mode information
            analytics_facts: Analytics facts
            initiatives: Selected initiatives
            vertical_name: Vertical name
            output_path: Output file path
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        self._add_title_slide(prs, company_name, vertical_name)
        
        # Slide 2: Executive Summary
        self._add_executive_summary(prs, mode_info, analytics_facts)
        
        # Slide 3: Key Metrics
        self._add_metrics_slide(prs, analytics_facts)
        
        # Slide 4: Initiatives Overview
        self._add_initiatives_overview(prs, initiatives)
        
        # Slides 5+: Individual initiatives (top 3)
        for init in initiatives[:3]:
            self._add_initiative_detail(prs, init)
        
        # Last slide: Next Steps
        self._add_next_steps(prs, mode_info)
        
        prs.save(output_path)
    
    def _add_title_slide(self, prs, company_name: str, vertical_name: str):
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = f"Diagnostic Report\n{company_name}"
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"{vertical_name} • {datetime.now().strftime('%B %Y')}"
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
    
    def _add_executive_summary(self, prs, mode_info: Dict[str, Any], analytics_facts: List[Dict[str, Any]]):
        """Add executive summary slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Operating mode
        p = tf.paragraphs[0]
        p.text = f"Operating Mode: {mode_info['mode']}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(12)
        
        # Confidence
        p = tf.add_paragraph()
        p.text = f"Confidence Level: {mode_info['confidence']}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        # Months
        p = tf.add_paragraph()
        p.text = f"Data Coverage: {mode_info['months_available']} months"
        p.font.size = Pt(16)
        p.space_after = Pt(20)
        
        # Key findings
        p = tf.add_paragraph()
        p.text = "Key Findings:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(8)
        
        for fact in analytics_facts[:3]:
            p = tf.add_paragraph()
            p.text = f"• {fact['label']}: {fact.get('value', 'N/A')} {fact.get('unit', '')}"
            p.font.size = Pt(14)
            p.level = 1
    
    def _add_metrics_slide(self, prs, analytics_facts: List[Dict[str, Any]]):
        """Add key metrics slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Key Metrics"
        
        # Create table
        rows = min(len(analytics_facts), 8) + 1
        cols = 3
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(0, 2).text = "Period"
        
        # Data rows
        for i, fact in enumerate(analytics_facts[:8]):
            table.cell(i + 1, 0).text = fact['label']
            value_str = f"{fact.get('value', fact.get('value_text', 'N/A'))}"
            if isinstance(fact.get('value'), (int, float)):
                if fact.get('unit') == 'currency':
                    value_str = f"${fact['value']:,.0f}"
                elif fact.get('unit') == 'percentage':
                    value_str = f"{fact['value']:.1f}%"
            table.cell(i + 1, 1).text = value_str
            table.cell(i + 1, 2).text = str(fact.get('period', ''))
    
    def _add_initiatives_overview(self, prs, initiatives: List[Dict[str, Any]]):
        """Add initiatives overview slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Recommended Initiatives"
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for init in initiatives:
            p = tf.add_paragraph() if init['rank'] > 1 else tf.paragraphs[0]
            p.text = f"{init['rank']}. {init['title']}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            impact_str = f"${init.get('impact_low', 0):,.0f} - ${init.get('impact_high', 0):,.0f}"
            p.text = f"   Estimated Impact: {impact_str}"
            p.font.size = Pt(14)
            p.level = 1
            p.space_after = Pt(12)
    
    def _add_initiative_detail(self, prs, initiative: Dict[str, Any]):
        """Add detailed initiative slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = f"{initiative['rank']}. {initiative['title']}"
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Category
        p = tf.paragraphs[0]
        p.text = f"Category: {initiative['category']}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        # Impact
        p = tf.add_paragraph()
        impact_str = f"${initiative.get('impact_low', 0):,.0f} - ${initiative.get('impact_high', 0):,.0f}"
        p.text = f"Estimated Annual Impact: {impact_str}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(16)
        
        # Explanation
        p = tf.add_paragraph()
        p.text = "Overview:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = initiative.get('explanation', initiative.get('description', ''))
        p.font.size = Pt(13)
        p.space_after = Pt(12)
        
        # Assumptions
        if initiative.get('assumptions'):
            p = tf.add_paragraph()
            p.text = "Key Assumptions:"
            p.font.size = Pt(14)
            p.font.bold = True
            p.space_after = Pt(6)
            
            for assumption in initiative['assumptions'][:3]:
                p = tf.add_paragraph()
                p.text = f"• {assumption}"
                p.font.size = Pt(12)
                p.level = 1
    
    def _add_next_steps(self, prs, mode_info: Dict[str, Any]):
        """Add next steps slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Next Steps"
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        steps = [
            "Review and validate recommended initiatives with stakeholders",
            "Prioritize 2-3 initiatives for immediate implementation",
            "Establish baseline metrics and tracking mechanisms",
            "Develop detailed implementation plans for selected initiatives",
            "Schedule follow-up diagnostic in 90 days to measure progress"
        ]
        
        for i, step in enumerate(steps):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"{i + 1}. {step}"
            p.font.size = Pt(16)
            p.space_after = Pt(16)
        
        # Confidence note
        p = tf.add_paragraph()
        p.text = f"\nAnalysis Confidence: {mode_info['confidence']} ({mode_info['mode']})"
        p.font.size = Pt(12)
        p.font.italic = True
